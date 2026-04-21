from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

OUT = 'llm_pricing_usage_report_presentation.pptx'
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

FONT = 'Noto Sans CJK SC'
TITLE_DARK = RGBColor(14, 30, 64)
ACCENT_BLUE = RGBColor(45, 110, 210)
ACCENT_ORANGE = RGBColor(243, 122, 41)
ACCENT_RED = RGBColor(176, 37, 37)
ACCENT_GREEN = RGBColor(38, 134, 83)
BG_LIGHT = RGBColor(246, 248, 252)
BG_DARK = RGBColor(12, 25, 52)
TEXT = RGBColor(38, 48, 62)
MUTED = RGBColor(98, 111, 130)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(214, 220, 230)


def set_cell_text(cell, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.CENTER):
    cell.text = ''
    tf = cell.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    font = run.font
    font.name = FONT
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                color=TEXT, align=PP_ALIGN.LEFT, fill=None, line=None,
                radius=False, margin=0.08, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for i, part in enumerate(str(text).split('\n')):
        if i == 0:
            run = p.add_run()
            para = p
        else:
            para = tf.add_paragraph()
            para.alignment = align
            run = para.add_run()
        run.text = part
        font = run.font
        font.name = FONT
        font.size = Pt(size)
        font.bold = bold
        font.color.rgb = color
    return shape


def add_title(slide, title, subtitle=None, dark=False):
    title_color = WHITE if dark else TITLE_DARK
    sub_color = RGBColor(214, 224, 243) if dark else MUTED
    add_textbox(slide, Inches(0.55), Inches(0.25), Inches(9.5), Inches(0.65), title,
                size=26, bold=True, color=title_color)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.83), Inches(11.8), Inches(0.35), subtitle,
                    size=11.5, color=sub_color)


def add_footer(slide, text, dark=False):
    color = RGBColor(210, 220, 235) if dark else MUTED
    add_textbox(slide, Inches(0.55), Inches(7.0), Inches(12.0), Inches(0.22), text,
                size=9.5, color=color)


def add_bullet_box(slide, left, top, width, height, title, bullets, accent=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = accent

    for bullet in bullets:
        p = tf.add_paragraph()
        p.level = 0
        p.bullet = True
        p.space_before = 0
        p.space_after = 0
        run = p.add_run()
        run.text = bullet
        run.font.name = FONT
        run.font.size = Pt(13)
        run.font.color.rgb = TEXT
    return shape


def add_table(slide, left, top, width, height, headers, rows, col_widths=None, header_fill=TITLE_DARK):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    row_h = int(height / (len(rows) + 1))
    for r in range(len(rows) + 1):
        table.rows[r].height = row_h
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        set_cell_text(cell, h, size=11.5, bold=True, color=WHITE)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else RGBColor(241, 245, 250)
            align = PP_ALIGN.LEFT if c == 0 or (isinstance(val, str) and len(val) > 18) else PP_ALIGN.CENTER
            set_cell_text(cell, val, size=10.5, bold=False, color=TEXT, align=align)
    return table


def add_bar_chart(slide, left, top, width, height, categories, values, title='', color=ACCENT_BLUE, max_scale=None):
    data = CategoryChartData()
    data.categories = categories
    data.add_series(title or '数值', values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(225, 230, 238)
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    if max_scale:
        chart.value_axis.maximum_scale = max_scale
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    series.format.line.color.rgb = color
    chart.chart_title.text_frame.text = title if title else ''
    if title:
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
    return chart


def add_stacked_chart(slide, left, top, width, height, categories, series_dict, colors):
    data = CategoryChartData()
    data.categories = categories
    for k, v in series_dict.items():
        data.add_series(k, v)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.font.size = Pt(10)
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(225, 230, 238)
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    for idx, s in enumerate(chart.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = colors[idx]
        s.format.line.color.rgb = colors[idx]
    return chart


def add_pie_chart(slide, left, top, width, height, categories, values, colors):
    data = CategoryChartData()
    data.categories = categories
    data.add_series('占比', values)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.font.size = Pt(10)
    points = chart.series[0].points
    for idx, pt in enumerate(points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[idx]
        pt.format.line.color.rgb = WHITE
    return chart


def add_connector_arrow(slide, left, top, width, height, text):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT_BLUE
    shp.line.color.rgb = ACCENT_BLUE
    tf = shp.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = FONT
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return shp


# Slide 1 Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)
# decorative bands
for i, c in enumerate([ACCENT_BLUE, ACCENT_GREEN, ACCENT_ORANGE]):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.8 + 0.25*i), Inches(0), Inches(0.28), Inches(7.5))
    band.fill.solid(); band.fill.fore_color.rgb = c; band.line.fill.background()
add_textbox(slide, Inches(0.75), Inches(1.0), Inches(7.1), Inches(1.9), '大模型定价、限制与\n研发团队月度用量评估', size=27, bold=True, color=WHITE)
add_textbox(slide, Inches(0.8), Inches(3.0), Inches(6.8), Inches(0.6), '围绕 GLM-5.1、Claude Opus、GPT-5 的选型与预算分析', size=16, color=RGBColor(216,223,236))
for idx, (t, c) in enumerate([('定价', ACCENT_ORANGE), ('限制', ACCENT_RED), ('用量', ACCENT_GREEN)]):
    add_textbox(slide, Inches(9.0), Inches(1.5 + idx*0.95), Inches(2.0), Inches(0.52), t,
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=c, radius=True, valign=MSO_ANCHOR.MIDDLE)
add_textbox(slide, Inches(0.8), Inches(6.55), Inches(5.5), Inches(0.35), '汇报人：Hermes Agent    日期：2026-04-15', size=10.5, color=RGBColor(188,201,220))

# Slide 2 conclusion
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '一页看结论', '预算关键不在需求分析，而在开发与 Review 的高频循环')
add_bullet_box(slide, Inches(0.55), Inches(1.25), Inches(3.6), Inches(1.55), '关键结论 1', [
    '20 人团队基线月消耗约 94.60M tokens',
    '默认预算建议先以基线场景作为编制口径'
], accent=ACCENT_BLUE)
add_bullet_box(slide, Inches(0.55), Inches(3.0), Inches(3.6), Inches(1.6), '关键结论 2', [
    '同等用量下：GLM-5.1 最低，GPT-5 次之，Claude Opus 4.1 最高',
    '全员默认高端模型会显著抬高月成本'
], accent=ACCENT_ORANGE)
add_bullet_box(slide, Inches(0.55), Inches(4.8), Inches(3.6), Inches(1.55), '关键结论 3', [
    '建议采用“低成本主力 + 高质量升级模型”组合',
    '复杂架构分析、关键 PR 二审再升级到 Opus'
], accent=ACCENT_GREEN)
add_bar_chart(slide, Inches(4.45), Inches(1.45), Inches(8.2), Inches(4.9),
              ['GLM-5.1', 'Claude Opus 4.1', 'GPT-5'], [143.11, 2871.00, 330.00], '基线场景月成本（美元）', color=ACCENT_BLUE, max_scale=3200)
add_footer(slide, '说明：Claude-opus-4.6 / GPT-5.4 非公开官方命名，报告分别按 Claude Opus 4.1 / GPT-5 对比。')

# Slide 3 scope
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '比较口径与边界说明', '先统一命名与价格口径，避免平台别名与官方 SKU 混淆')
# flow
boxes = [
    ('原始型号', 'Claude-opus-4.6\nGPT-5.4\nGLM-5.1', RGBColor(233, 238, 247), TITLE_DARK),
    ('纠偏口径', 'Claude Opus 4.1\nGPT-5\nGLM-5.1', RGBColor(227, 241, 234), ACCENT_GREEN),
    ('价格来源', '官方优先\nGLM 价格采用 OpenRouter\n仅作横向参考', RGBColor(255, 242, 230), ACCENT_ORANGE),
]
lefts = [0.7, 4.4, 8.1]
for (title, body, fill, accent), l in zip(boxes, lefts):
    add_textbox(slide, Inches(l), Inches(1.9), Inches(3.0), Inches(2.5), f'{title}\n{body}', size=15,
                bold=False, color=TEXT, fill=fill, line=accent, radius=True)
add_connector_arrow(slide, Inches(3.8), Inches(2.8), Inches(0.45), Inches(0.5), '→')
add_connector_arrow(slide, Inches(7.5), Inches(2.8), Inches(0.45), Inches(0.5), '→')
add_bullet_box(slide, Inches(0.8), Inches(4.85), Inches(5.7), Inches(1.3), '边界与风险', [
    '官方信息抓取不稳定时，采用公开聚合平台页面做辅助参考',
    '不同路由平台、账号等级、采购方式，实际账单可能不同'
], accent=ACCENT_RED)
add_bullet_box(slide, Inches(6.75), Inches(4.85), Inches(5.7), Inches(1.3), '本页核心目的', [
    '防止后续讨论时把“平台别名”与“官方型号”混为一谈',
    '确保后续预算测算按同一比较口径展开'
], accent=ACCENT_BLUE)
add_footer(slide, 'GLM-5.1：官方可确认型号存在；公开价格结构未稳定提取，因此采用 OpenRouter 转售口径做横向参考。')

# Slide 4 pricing table
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '主流模型价格与限制横向对比', '重点比较输入/输出价格、上下文窗口与限制口径')
headers = ['对比型号', '输入价格', '输出价格', '缓存价格', '上下文窗口', '最大输出', '限制口径']
rows = [
    ['GLM-5.1', '$0.95 / 1M', '$3.15 / 1M', 'Read $0.475 / 1M', '202,752', '65,535', '公开统一限流未明确抓取'],
    ['Claude Opus 4.1', '$15 / 1M', '$75 / 1M', 'Write $18.75 / 1M\nRead $1.50 / 1M', '200K', '32K', 'tier-based'],
    ['GPT-5', '$1.25 / 1M', '$10 / 1M', 'Cached input\n$0.125 / 1M', '400K', '128K', 'tier-based'],
]
add_table(slide, Inches(0.45), Inches(1.4), Inches(12.4), Inches(3.2), headers, rows,
          [Inches(1.6), Inches(1.3), Inches(1.3), Inches(2.2), Inches(1.3), Inches(1.1), Inches(2.6)])
add_textbox(slide, Inches(0.6), Inches(5.0), Inches(3.8), Inches(1.0), '直观认知\nClaude 最贵，GPT 次之，GLM 当前参考口径最低',
            size=16, bold=True, color=TITLE_DARK, fill=RGBColor(235, 242, 255), line=ACCENT_BLUE, radius=True)
add_textbox(slide, Inches(4.7), Inches(5.0), Inches(4.0), Inches(1.0), '注意事项\n缓存价格不能直接等同于基础 input 单价',
            size=14, color=TEXT, fill=RGBColor(255, 245, 236), line=ACCENT_ORANGE, radius=True)
add_textbox(slide, Inches(9.0), Inches(5.0), Inches(3.3), Inches(1.0), '注释\nGLM 价格为转售参考口径',
            size=14, color=TEXT, fill=RGBColor(245, 238, 238), line=ACCENT_RED, radius=True)
add_footer(slide, '表格用于形成管理层直观认知，不展开 RPM/TPM 细节；实际限额仍以账号等级与采购平台配置为准。')

# Slide 5 token source
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, 'Token 消耗的主要来源在哪里', '研发场景不是一次性问答，而是天然形成多轮循环')
steps = [
    ('需求生成', '单次对话较长\n但频次相对低', RGBColor(226, 236, 250), ACCENT_BLUE, 1.15),
    ('开发', '高频\n携带代码上下文 / 日志', RGBColor(255, 241, 230), ACCENT_ORANGE, 1.45),
    ('Review', '多轮反复\ndiff / 修改结果持续回放', RGBColor(252, 228, 228), ACCENT_RED, 1.55),
    ('测试', '环境排障与测例编写\n需要多轮追问', RGBColor(231, 244, 236), ACCENT_GREEN, 1.2),
    ('回流修复', '问题闭环后\n再次进入开发与 Review', RGBColor(238, 241, 246), TITLE_DARK, 1.1),
]
cur = 0.6
for idx, (name, desc, fill, line, h) in enumerate(steps):
    add_textbox(slide, Inches(cur), Inches(2.3), Inches(2.2), Inches(h), f'{name}\n{desc}', size=14,
                color=TEXT, fill=fill, line=line, radius=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if idx < len(steps)-1:
        add_connector_arrow(slide, Inches(cur+2.25), Inches(2.72), Inches(0.38), Inches(0.45), '→')
    cur += 2.55
add_bullet_box(slide, Inches(0.9), Inches(5.2), Inches(5.5), Inches(1.15), '为什么开发 / Review 最耗 token', [
    '交互轮次最多，且每轮都会重复携带上下文',
    '代码、日志、diff 与修改结果叠加放大输入 token'
], accent=ACCENT_ORANGE)
add_bullet_box(slide, Inches(6.7), Inches(5.2), Inches(5.5), Inches(1.15), '管理提示', [
    '如果要控预算，优先优化 Review loop 与上下文复用策略',
    '需求生成通常不是成本大头'
], accent=ACCENT_GREEN)

# Slide 6 usage scenarios
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '20 人团队三档月用量评估', '建议把“基线场景”作为预算编制默认口径')
headers = ['场景', '月输入 tokens', '月输出 tokens', '月总 tokens']
rows = [
    ['保守', '47.96M', '16.28M', '64.24M'],
    ['基线', '70.40M', '24.20M', '94.60M'],
    ['激进', '127.60M', '44.00M', '171.60M'],
]
add_table(slide, Inches(0.65), Inches(1.55), Inches(4.4), Inches(2.2), headers, rows,
          [Inches(0.95), Inches(1.1), Inches(1.1), Inches(1.25)])
add_stacked_chart(slide, Inches(5.3), Inches(1.4), Inches(7.0), Inches(4.2), ['保守', '基线', '激进'],
                  {'输入': [47.96, 70.40, 127.60], '输出': [16.28, 24.20, 44.00]},
                  [ACCENT_BLUE, ACCENT_ORANGE])
add_textbox(slide, Inches(0.95), Inches(4.1), Inches(3.7), Inches(1.1), '基线场景\n94.60M tokens / 月',
            size=18, bold=True, color=WHITE, fill=ACCENT_BLUE, line=ACCENT_BLUE, radius=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_bullet_box(slide, Inches(0.7), Inches(5.45), Inches(4.6), Inches(0.8), '场景解释', [
    '保守：刚开始试点；激进：深度 AI-native；基线：已进入日常研发流程'
], accent=ACCENT_BLUE)
add_footer(slide, '堆叠柱展示输入 / 输出拆分，便于后续直接映射不同模型的 input/output 计费单价。')

# Slide 7 baseline split
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '基线场景下，成本主要消耗在哪些环节', '基线总量 94.60M tokens，其中代码开发占比最高')
headers = ['环节', '月输入', '月输出', '月总量']
rows = [
    ['需求生成', '11.88M', '4.84M', '16.72M'],
    ['代码开发', '38.72M', '13.20M', '51.92M'],
    ['测试验证', '19.80M', '6.16M', '25.96M'],
]
add_table(slide, Inches(0.65), Inches(1.55), Inches(4.5), Inches(2.2), headers, rows,
          [Inches(1.0), Inches(1.1), Inches(1.1), Inches(1.3)])
add_pie_chart(slide, Inches(5.45), Inches(1.45), Inches(3.25), Inches(3.4),
              ['需求生成', '代码开发', '测试验证'], [16.72, 51.92, 25.96], [ACCENT_BLUE, ACCENT_ORANGE, ACCENT_GREEN])
add_bar_chart(slide, Inches(8.75), Inches(1.55), Inches(3.6), Inches(3.2), ['需求', '开发', '测试'], [16.72, 51.92, 25.96], '基线月总量（M）', color=ACCENT_ORANGE, max_scale=60)
add_bullet_box(slide, Inches(0.9), Inches(5.2), Inches(11.3), Inches(0.95), '预算控制提示', [
    '若后续要控预算，优先优化开发与 Review 的调用策略，而不是先压缩需求分析环节'
], accent=ACCENT_ORANGE)

# Slide 8 cost compare
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '同样用量下，不同模型月成本差异明显', '基线场景成本未计缓存命中、企业折扣、Batch 与转售差价')
headers = ['模型', '基线月成本']
rows = [['GLM-5.1', '$143.11'], ['Claude Opus 4.1', '$2,871.00'], ['GPT-5', '$330.00']]
add_table(slide, Inches(0.8), Inches(1.75), Inches(3.2), Inches(1.9), headers, rows, [Inches(1.6), Inches(1.4)])
add_bar_chart(slide, Inches(4.3), Inches(1.4), Inches(8.0), Inches(4.7), ['GLM-5.1', 'Claude Opus 4.1', 'GPT-5'], [143.11, 2871.0, 330.0], '基线月成本对比（美元）', color=ACCENT_RED, max_scale=3200)
add_textbox(slide, Inches(0.82), Inches(4.2), Inches(3.15), Inches(1.25), '解读\nClaude 更适合高价值、高难度任务，\n不宜作为全员日常默认底座',
            size=14, color=TEXT, fill=RGBColor(255, 240, 240), line=ACCENT_RED, radius=True)
add_footer(slide, '最适合回答“如果全员默认都用高端模型，会贵多少”。')

# Slide 9 layered strategy
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '推荐采用“分层模型策略”', '真正高性价比通常不是“最强单模型”，而是“主力模型 + 专家模型”组合')
add_textbox(slide, Inches(0.75), Inches(1.55), Inches(4.8), Inches(3.8), '默认流量\n\n推荐模型：GPT-5 或 GLM-5.1\n\n适用：\n• 日常开发\n• 常规 Review\n• 测试辅助\n\n目标：\n覆盖高频、标准化、成本敏感任务',
            size=16, color=TEXT, fill=RGBColor(235, 242, 255), line=ACCENT_BLUE, radius=True)
add_textbox(slide, Inches(7.75), Inches(1.55), Inches(4.8), Inches(3.8), '升级流量\n\n推荐模型：Claude Opus 4.1\n\n适用：\n• 复杂架构分析\n• 高风险 PR 二次 review\n• 疑难性能问题定位\n\n目标：\n在关键复杂任务上保留高质量能力',
            size=16, color=TEXT, fill=RGBColor(233, 244, 237), line=ACCENT_GREEN, radius=True)
add_connector_arrow(slide, Inches(5.95), Inches(2.7), Inches(1.15), Inches(0.75), '升级条件')
add_bullet_box(slide, Inches(1.2), Inches(5.65), Inches(11.0), Inches(0.6), '策略收益', [
    '既保留复杂任务质量，又避免日常全量流量落在高价模型上导致成本失控'
], accent=ACCENT_GREEN)

# Slide 10 next steps
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)
add_title(slide, '预算、治理与落地建议', '建议先试运行 1 个月，再用真实调用日志校准预算模型')
items = [
    ('1', '以“基线场景 94.60M tokens/月”作为预算测算起点', ACCENT_BLUE),
    ('2', '默认启用缓存、上下文复用、模板化 prompt', ACCENT_GREEN),
    ('3', '对 Review loop 设置轮次阈值，避免无上限反复对话', ACCENT_ORANGE),
    ('4', '将环境排障、测例编写优先落到低成本模型', ACCENT_BLUE),
    ('5', '对关键任务保留高端模型升级通道', ACCENT_RED),
]
for idx, (n, txt, c) in enumerate(items):
    y = 1.5 + idx*0.88
    add_textbox(slide, Inches(0.95), Inches(y), Inches(0.55), Inches(0.45), n,
                size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, fill=c, line=c, radius=True, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(1.7), Inches(y-0.02), Inches(10.5), Inches(0.5), txt,
                size=16, color=TEXT, fill=WHITE, line=BORDER, radius=True, valign=MSO_ANCHOR.MIDDLE)
add_textbox(slide, Inches(0.95), Inches(6.25), Inches(11.25), Inches(0.7),
            '一句话收尾：建议采用低成本主力模型承接日常流量，以高质量模型承接关键复杂任务，在保证效果的同时控制团队级成本。',
            size=15, bold=True, color=WHITE, fill=TITLE_DARK, line=TITLE_DARK, radius=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(OUT)
