from pathlib import Path

from pptx import Presentation

from ppt.generate_llm_pricing_report import build_presentation


def _collect_text(prs: Presentation) -> str:
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def test_build_presentation_creates_expected_slide_count(tmp_path: Path) -> None:
    output = tmp_path / "report.pptx"
    build_presentation(output)

    assert output.exists()
    prs = Presentation(str(output))
    assert len(prs.slides) == 12


def test_build_presentation_contains_expected_titles(tmp_path: Path) -> None:
    output = tmp_path / "report.pptx"
    build_presentation(output)

    prs = Presentation(str(output))
    titles = [
        slide.shapes.title.text
        for slide in prs.slides
        if slide.shapes.title is not None
    ]

    assert "执行摘要" in titles
    assert "三家模型价格总览" in titles
    assert "20 人团队月度 Token 基线模型" in titles
    assert "基线月费测算" in titles


def test_build_presentation_contains_key_metrics(tmp_path: Path) -> None:
    output = tmp_path / "report.pptx"
    build_presentation(output)

    prs = Presentation(str(output))
    merged = _collect_text(prs)

    assert "618.2M" in merged
    assert "5,102 元 ~ 6,252 元/月" in merged
    assert "$2,828 /月" in merged
