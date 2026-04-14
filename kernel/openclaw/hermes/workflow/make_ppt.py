from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

out = Path('ai-assisted-kernel-development-report.pptx')
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(0x16,0x21,0x3E)
TEAL = RGBColor(0x0F,0x76,0x73)
CYAN = RGBColor(0x22,0xB8,0xCF)
GOLD = RGBColor(0xF5,0xB7,0x00)
RED = RGBColor(0xD9,0x48,0x41)
GREEN = RGBColor(0x2B,0x8A,0x3E)
LIGHT = RGBColor(0xF7,0xF9,0xFC)
TEXT = RGBColor(0x20,0x2B,0x38)
MUTED = RGBColor(0x5C,0x6B,0x7A)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BORDER = RGBColor(0xC8,0xD2,0xDC)

prs.core_properties.author = 'Hermes'
prs.core_properties.title = 'AI辅助内核开发多智能体工作流汇报'
prs.core_properties.subject = '探索-规划-生成-审核-调试 多智能体方案'

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def textbox(slide, x, y, w, h, text='', size=20, color=TEXT, bold=False, font='Microsoft YaHei', align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    p.alignment = align
    return tb

def add_bullets(slide, x, y, w, h, items, size=18, color=TEXT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            level, txt = item
        else:
            level, txt = 0, item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = txt
        p.level = level
        p.font.name = 'Microsoft YaHei'
        p.font.size = Pt(size - level)
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = Pt(7)
    return tb

def card(slide, x, y, w, h, title, body_lines, accent=TEAL, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x+0.25, y+0.16, w-0.35, 0.48, title, size=18, color=TEXT, bold=True)
    add_bullets(slide, x+0.22, y+0.72, w-0.34, h-0.85, body_lines, size=14, color=MUTED)

def pill(slide, x, y, w, h, text, fillc, txtc=WHITE, size=16):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fillc
    s.line.fill.background()
    textbox(slide, x, y+0.03, w, h-0.02, text, size=size, color=txtc, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)

def arrow(slide, x1, y1, x2, y2, color=CYAN, width=2.5):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c

# 1
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, NAVY)
for x,y,w,h,c in [(0.6,0.7,1.2,0.18,TEAL),(1.95,0.7,0.55,0.18,CYAN),(2.65,0.7,0.75,0.18,GOLD)]:
    s=slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)); s.fill.solid(); s.fill.fore_color.rgb=c; s.line.fill.background()
textbox(slide, 0.85, 1.25, 8.9, 1.25, 'AI辅助内核开发\n多智能体工作流汇报', size=28, color=WHITE, bold=True)
pill(slide, 0.88, 2.65, 2.5, 0.46, '探索 → 规划 → 生成', TEAL, size=16)
pill(slide, 3.55, 2.65, 2.15, 0.46, '审核 ↔ 调试', CYAN, txtc=NAVY, size=16)
pill(slide, 5.92, 2.65, 2.2, 0.46, 'Patch Ready', GOLD, txtc=NAVY, size=16)
textbox(slide, 0.9, 3.45, 8.8, 1.2, '面向 Linux 内核 / RISC-V / KVM 等高上下文、高验证成本场景\n将“问题发现到补丁准备”组织为可审计、可迭代、带人工闸门的工程系统', size=18, color=RGBColor(0xD6,0xE4,0xF0))
shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(1.15), Inches(2.7), Inches(4.9))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x20,0x2E,0x55); shape.line.color.rgb = RGBColor(0x3B,0x4A,0x72)
for i,(txt,c) in enumerate([('探索',TEAL),('规划',CYAN),('生成',GOLD),('审核',RGBColor(0xF0,0x6E,0x6E)),('调试',GREEN)]):
    yy = 1.55 + i*0.78
    pill(slide, 10.0, yy, 1.7, 0.42, txt, c, txtc=(NAVY if c in [CYAN,GOLD] else WHITE), size=15)
    if i<4: arrow(slide, 10.85, yy+0.42, 10.85, yy+0.72, color=WHITE, width=1.8)
textbox(slide, 0.9, 6.55, 5, 0.35, '汇报内容：方案概览、架构、迭代闭环、示例、落地建议', size=11, color=RGBColor(0x9F,0xB3,0xC8))

# 2
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
textbox(slide, 0.7, 0.45, 5.5, 0.5, '为什么需要这套工作流', size=28, color=NAVY, bold=True)
textbox(slide, 0.72, 0.98, 9.5, 0.35, '核心判断：内核开发不是“让一个模型写代码”，而是把高风险认知活动拆成可验证的职责链。', size=15, color=MUTED)
card(slide, 0.75, 1.45, 4.0, 2.0, '当前挑战', ['内核语义强依赖架构背景、邮件列表上下文和历史设计取舍', '一次性大任务容易导致 token 开销、实现偏航和审查风险同步放大', '“能编过”不等于“能被上游接受”'], accent=RED)
card(slide, 4.95, 1.45, 3.55, 2.0, '设计原则', ['小问题切分，单轮只处理一个窄议题', '状态与证据外置为工件，避免上下文膨胀', '关键节点保留人工 Gate'], accent=TEAL)
card(slide, 8.7, 1.45, 3.85, 2.0, '目标结果', ['建立 discover → plan → generate → review → debug → patch 的闭环', '形成可审计、可重试、可复盘的研发流水线', '把人工时间集中在高价值判断上'], accent=GOLD)
for x,w,num,label,c in [(0.85,2.7,'5','核心智能体层级',NAVY),(3.9,2.7,'3','人工闸门',TEAL),(6.95,2.7,'2+','生成/审核迭代轮次',CYAN),(10.0,2.2,'1','issue / branch / target per run',GOLD)]:
    s=slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(4.15), Inches(w), Inches(1.55))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE; s.line.color.rgb = BORDER
    textbox(slide, x, 4.42, w, 0.45, num, size=26, color=c, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x+0.1, 4.95, w-0.2, 0.35, '每轮仅 1 个 issue/分支/目标' if num == '1' else label, size=13, color=MUTED, align=PP_ALIGN.CENTER)
textbox(slide, 0.85, 6.15, 11.8, 0.5, '结论：AI 更适合承担“探索、方案生成、验证、归因、修复”的模块化工作，而不是跳过流程直接投递上游。', size=16, color=TEXT)

# 3
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
textbox(slide, 0.7, 0.45, 6.5, 0.5, '多智能体总体架构', size=28, color=NAVY, bold=True)
textbox(slide, 0.72, 0.98, 11.2, 0.35, '控制中枢负责编排；探索、规划、生成、审核、调试五层分工协作；人类只在高风险节点审批。', size=15, color=MUTED)
s=slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.45), Inches(1.55), Inches(2.4), Inches(0.85)); s.fill.solid(); s.fill.fore_color.rgb=NAVY; s.line.fill.background()
textbox(slide,5.45,1.74,2.4,0.35,'Orchestrator',size=21,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
pill(slide, 5.8, 0.92, 1.7, 0.42, 'Human Gates', RED, size=15)
arrow(slide, 6.65, 1.34, 6.65, 1.55, color=RED)
card(slide, 0.8, 1.65, 3.8, 1.55, '探索层', ['Scout-Code：找代码差距', 'Scout-History：找 lore/历史证据', 'Triage：去重、评分、过滤伪问题'], accent=TEAL)
card(slide, 0.8, 3.5, 3.8, 1.55, '规划层', ['Planner：file-level 最小方案', 'Test-Designer：测试矩阵', 'Risk-Reviewer：ABI/UAPI/回滚风险'], accent=CYAN)
card(slide, 8.7, 1.4, 3.8, 1.3, '生成层', ['Implementer-A：主实现', 'Implementer-B：备选实现/worktree'], accent=GOLD)
card(slide, 8.7, 2.95, 3.8, 1.5, '审核层', ['Spec-Review：是否满足设计目标', 'Code-Review：质量/边界/并发', 'Upstream-Review：patch 粒度与提交叙事'], accent=RED)
card(slide, 8.7, 4.75, 3.8, 1.45, '调试层', ['Failure-Analyzer：根因分析', 'Fix-Agent：定向修复', 'Regression-Guard：防回归'], accent=GREEN)
arrow(slide,4.6,2.4,5.45,1.98,color=TEAL); arrow(slide,4.6,4.25,5.45,1.98,color=CYAN)
arrow(slide,7.85,1.98,8.7,2.05,color=GOLD); arrow(slide,7.85,1.98,8.7,3.6,color=RED); arrow(slide,7.85,1.98,8.7,5.45,color=GREEN)
textbox(slide, 0.95, 6.45, 11.6, 0.45, '关键约束：同一 issue 在生成阶段只能有一个主写分支；审核可并行；调试必须串行回到审核。', size=14, color=TEXT)

# 4
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
textbox(slide, 0.7, 0.45, 8.5, 0.5, '五阶段主线：探索 → 规划 → 生成 → 审核 → 调试', size=26, color=NAVY, bold=True)
xs=[0.8,3.0,5.2,7.4,9.6]
labels=[('探索','发现真实问题',TEAL),('规划','形成最小方案',CYAN),('生成','产出代码/日志',GOLD),('审核','独立质量判定',RED),('调试','失败闭环与回归防护',GREEN)]
for i,(x,(a,b,c)) in enumerate(zip(xs,labels)):
    pill(slide,x,1.15,1.8,0.48,a,c,txtc=(NAVY if c in [CYAN,GOLD] else WHITE),size=17)
    textbox(slide,x-0.1,1.7,2.0,0.4,b,size=12,color=MUTED,align=PP_ALIGN.CENTER)
    if i<4: arrow(slide,x+1.85,1.39,xs[i+1]-0.12,1.39,color=BORDER,width=2)
stage_data=[
('探索',0.7,['输入：源码、架构目录、lore、历史提交','输出：代码证据、历史证据、gap_registry','DoD：问题真实、范围清晰、可转为窄 issue'],TEAL),
('规划',3.2,['输入：gap 条目与证据 bundle','输出：design / test matrix / risk','DoD：文件级路径清晰、可执行、可回滚'],CYAN),
('生成',5.7,['输入：方案与测试矩阵','输出：代码修改、构建日志、测试日志、draft patch','DoD：最小实现可被 review 消费'],GOLD),
('审核',8.2,['输入：实现结果与日志工件','输出：spec/code/upstream 三类评审','DoD：PASS / REVISE / DEBUG / BLOCKED'],RED),
('调试',10.7,['输入：失败日志与审查意见','输出：根因报告、修复摘要、回归结果','DoD：回到审核而不是直接放行'],GREEN)
]
for title,x,items,c in stage_data:
    shape=slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(2.0), Inches(3.75))
    shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT; shape.line.color.rgb = BORDER
    bar=slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(2.3), Inches(2.0), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb=c; bar.line.fill.background()
    textbox(slide,x+0.12,2.48,1.76,0.35,title,size=18,color=TEXT,bold=True,align=PP_ALIGN.CENTER)
    add_bullets(slide,x+0.12,2.92,1.72,2.9,items,size=12,color=MUTED)

# 5
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
textbox(slide, 0.7, 0.45, 7.2, 0.5, '核心机制：生成与审核的多轮迭代', size=28, color=NAVY, bold=True)
textbox(slide, 0.72, 0.98, 11.4, 0.4, '原则：生成不是一次成稿，审核也不是终点判决；两者之间必须允许“修订生成 ↔ 再审核”的多轮来回。', size=15, color=MUTED)
pill(slide,1.0,2.0,2.2,0.56,'Generate R1',GOLD,txtc=NAVY,size=18)
pill(slide,4.0,2.0,2.2,0.56,'Review R1',RED,size=18)
pill(slide,7.0,2.0,2.2,0.56,'Debug R1',GREEN,size=18)
pill(slide,10.0,2.0,2.2,0.56,'Review R2',RED,size=18)
arrow(slide,3.25,2.28,4.0,2.28,color=GOLD)
arrow(slide,6.25,2.28,7.0,2.28,color=RED)
arrow(slide,9.25,2.28,10.0,2.28,color=GREEN)
arrow(slide,10.95,2.58,10.95,3.45,color=CYAN)
arrow(slide,10.95,3.45,2.1,3.45,color=CYAN)
arrow(slide,2.1,3.45,2.1,2.58,color=CYAN)
textbox(slide,2.65,3.18,7.8,0.3,'若仍需修订，则进入 Generate R2 / R3，直到 PASS 或触发人工 Gate',size=13,color=CYAN,align=PP_ALIGN.CENTER)
card(slide, 0.95, 4.25, 3.8, 1.8, '每轮必须记录什么', ['round / issue_id / 输入工件路径', '本轮目标、修改文件、构建/测试结果', 'review decision 与 next action'], accent=TEAL)
card(slide, 4.8, 4.25, 3.8, 1.8, '停止条件', ['连续 3 轮无实质收敛', '同一错误反复出现且根因不稳定', '出现架构语义争议或回归持续增加'], accent=RED)
card(slide, 8.65, 4.25, 3.8, 1.8, '价值', ['降低“一次性过度设计”风险', '把失败从黑盒变成结构化反馈', '让人工只介入真正难判的地方'], accent=GOLD)

# 6
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
textbox(slide, 0.7, 0.45, 8.8, 0.5, '示例：demo-riscv-kvm-mmu-kconfig', size=27, color=NAVY, bold=True)
textbox(slide, 0.72, 0.98, 11.4, 0.45, '演示问题：让 CONFIG_RISCV_KVM 的 MMU 依赖在 Kconfig dependency 与 help text 中都显式表达。', size=15, color=MUTED)
pill(slide, 0.95, 1.65, 2.4, 0.48, '1. Issue 定义', NAVY, size=16)
pill(slide, 3.65, 1.65, 2.55, 0.48, '2. Round 1 生成', GOLD, txtc=NAVY, size=16)
pill(slide, 6.55, 1.65, 2.55, 0.48, '3. 审核驳回 + 调试', RED, size=16)
pill(slide, 9.45, 1.65, 2.7, 0.48, '4. Round 2 通过', GREEN, size=16)
for x1,x2,c in [(3.35,3.65,NAVY),(6.25,6.55,GOLD),(9.15,9.45,RED)]: arrow(slide,x1,1.89,x2,1.89,color=c)
card(slide,0.8,2.45,2.7,3.35,'Issue Brief',['低风险 Kconfig 一致性修复','不涉及 UAPI/ABI','适合演示：生成→审核→调试→再审核'],accent=NAVY)
card(slide,3.6,2.45,2.8,3.35,'Round 1 实现',['Implementer 只修了 dependency','构建/配置检查未失败','但未覆盖 help text，规格未完成'],accent=GOLD)
card(slide,6.5,2.45,2.9,3.35,'Spec / Debug',['Spec-Review 决策：ENTER_DEBUG','Failure-Analyzer 认定为 spec gap','Fix-Agent 建议仅继续修改 arch/riscv/kvm/Kconfig'],accent=RED)
card(slide,9.5,2.45,2.9,3.35,'Round 2 结果',['补充 help text 的 MMU 说明','保持单文件、低风险、可回滚','Spec Review Round 2：PASS，进入 patch-ready'],accent=GREEN)
textbox(slide,0.9,6.3,11.4,0.42,'这个示例证明：调试不一定来自编译失败，也可以来自“规格审查指出需求未闭环”。',size=15,color=TEXT)

# 7
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
textbox(slide, 0.7, 0.45, 8.2, 0.5, '工件体系、质量闸门与调度建议', size=27, color=NAVY, bold=True)
card(slide,0.75,1.3,5.2,4.6,'关键工件目录',['artifacts/<issue>/discover：代码/历史证据','plans：design / test-matrix / risk','run_history + logs：每轮执行与验证记录','review：spec/code/upstream 审查结论','debug：failure-analysis / fix / regression','patches：draft patch / cover letter / checkpatch'],accent=TEAL)
card(slide,6.2,1.3,2.9,2.15,'三个人工 Gate',['Gate-1：问题是否真实','Gate-2：方案是否最小可行','Gate-3：patch/cover/收件人是否合规'],accent=RED)
card(slide,9.35,1.3,3.0,2.15,'自动质量闸门',['make ARCH=riscv defconfig / build','kselftest / kunit / QEMU 板卡验证','checkpatch.pl / get_maintainer.pl'],accent=GOLD)
card(slide,6.2,3.8,6.15,2.1,'调度建议',['discovery-job：定时发现 gap', 'planning-job：issue_created 后生成方案', 'generation-job：design_ready 后产出代码', 'review-job：独立并行审查', 'debug-job：review 或 test failed 后进入修复'],accent=CYAN)
textbox(slide,0.85,6.3,11.6,0.42,'建议：每个 issue 独立目录、独立状态机、独立 worktree，所有 Agent 只读取必要摘要与路径，不吞整段历史上下文。',size=14,color=TEXT)

# 8
slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, NAVY)
textbox(slide,0.75,0.55,7.5,0.55,'落地路径与汇报结论',size=28,color=WHITE,bold=True)
card(slide,0.85,1.45,3.6,3.6,'MVP 落地建议',['先只跑 1 个高置信度窄 issue', '至少经历 2 轮生成-审核迭代', '至少经历 1 次调试-再审核闭环', '生成 patch-ready 材料，但保留人工最终发信'],accent=TEAL,fill=RGBColor(0xF5,0xF8,0xFC))
card(slide,4.85,1.45,3.6,3.6,'下一步工程化',['补 workflow.yaml / 目录模板', '固化 Claude Code / Codex 角色 prompts', '脚本化 checkpatch / 日志归档 / 状态回写', '加入预算控制与失败阈值'],accent=CYAN,fill=RGBColor(0xF5,0xF8,0xFC))
card(slide,8.85,1.45,3.6,3.6,'最终结论',['多智能体价值不在“多开模型”，而在职责解耦与相互制衡', '生成与审核必须多轮迭代', '调试必须回流审核，防止“修完即过”', '人工应只在高价值判断点介入'],accent=GOLD,fill=RGBColor(0xF5,0xF8,0xFC))
pill(slide,4.25,5.7,4.9,0.58,'半自动 · 强审计 · 可迭代 · 面向内核贡献',GOLD,txtc=NAVY,size=20)
textbox(slide,0.8,6.7,6.0,0.3,'文件：ai-assisted-kernel-development-report.pptx',size=11,color=RGBColor(0xB7,0xC4,0xD6))

prs.save(out)
print(out)
print('slides=', len(prs.slides))
