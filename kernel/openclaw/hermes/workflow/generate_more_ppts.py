from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

NAVY = RGBColor(0x16,0x21,0x3E)
TEAL = RGBColor(0x0F,0x76,0x73)
CYAN = RGBColor(0x22,0xB8,0xCF)
GOLD = RGBColor(0xF5,0xB7,0x00)
RED = RGBColor(0xD9,0x48,0x41)
GREEN = RGBColor(0x2B,0x8A,0x3E)
LIGHT = RGBColor(0xF7,0xF9,0xFC)
TEXT = RGBColor(0x20,0x2B,0x38)
MUTED = RGBColor(0x5C,0x6B,0x7A)
WHITE = RGBColor(0xFF,0xFF,0xFF)
BORDER = RGBColor(0xC8,0xD2,0xDC)
PALE = RGBColor(0xED, 0xF4, 0xF8)
SOFT = RGBColor(0xF3, 0xF6, 0xFB)

W = 13.333
H = 7.5


def new_prs(title, subject):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.author = 'Hermes'
    prs.core_properties.title = title
    prs.core_properties.subject = subject
    return prs


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def textbox(slide, x, y, w, h, text='', size=18, color=TEXT, bold=False,
            font='Microsoft YaHei', align=PP_ALIGN.LEFT,
            valign=MSO_VERTICAL_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    p.alignment = align
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=TEXT, font='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in items:
        level, txt = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = txt
        p.level = level
        p.font.name = font
        p.font.size = Pt(size - level)
        p.font.color.rgb = color
        p.bullet = True
        p.space_after = Pt(5)
    return tb


def card(slide, x, y, w, h, title, items, accent=TEAL, fill=WHITE, title_size=18, body_size=13, title_color=TEXT, body_color=MUTED):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = BORDER
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.12), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    textbox(slide, x+0.22, y+0.15, w-0.3, 0.38, title, size=title_size, color=title_color, bold=True)
    bullets(slide, x+0.2, y+0.62, w-0.3, h-0.75, items, size=body_size, color=body_color)


def pill(slide, x, y, w, h, text, fillc, txtc=WHITE, size=15):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fillc
    s.line.fill.background()
    textbox(slide, x, y+0.02, w, h-0.04, text, size=size, color=txtc, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)


def arrow(slide, x1, y1, x2, y2, color=CYAN, width=2.0):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def codebox(slide, x, y, w, h, title, code, accent=NAVY, fill=RGBColor(0x11,0x1A,0x2E), text_color=RGBColor(0xE9,0xF1,0xFA), size=11):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    outer.fill.solid(); outer.fill.fore_color.rgb = fill
    outer.line.color.rgb = accent
    pill(slide, x+0.18, y+0.12, min(2.2, w-0.36), 0.34, title, accent, size=13)
    textbox(slide, x+0.18, y+0.58, w-0.36, h-0.72, code, size=size, color=text_color, font='Consolas')


def stat_card(slide, x, y, w, h, value, label, accent):
    s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = BORDER
    textbox(slide, x, y+0.14, w, 0.42, value, size=24, color=accent, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, x+0.08, y+0.62, w-0.16, 0.34, label, size=12, color=MUTED, align=PP_ALIGN.CENTER)


def read(path):
    return Path(path).read_text(encoding='utf-8')


def snippet(path, start=0, end=12):
    lines = read(path).splitlines()[start:end]
    return '\n'.join(lines)


def make_tech_review():
    prs = new_prs('AI辅助内核开发多智能体工作流：技术评审细节版', 'technical review version')

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, NAVY)
    textbox(slide, 0.8, 1.0, 8.6, 1.2, 'AI辅助内核开发\n多智能体工作流（技术评审细节版）', size=28, color=WHITE, bold=True)
    textbox(slide, 0.82, 2.45, 8.8, 0.9, '面向技术方案评审：强调状态机、工件契约、失败闭环、质量闸门、预算控制与 patch-ready 条件。', size=17, color=RGBColor(0xD6,0xE4,0xF0))
    pill(slide, 0.82, 3.55, 2.0, 0.44, 'Artifact-Driven', TEAL)
    pill(slide, 3.0, 3.55, 2.1, 0.44, 'Review-Oriented', CYAN, txtc=NAVY)
    pill(slide, 5.28, 3.55, 1.9, 0.44, 'Patch-Ready', GOLD, txtc=NAVY)
    card(slide, 9.4, 1.1, 2.9, 4.9, '评审关注点', ['职责边界是否清晰', '状态回写是否可追踪', '生成/审核能否多轮收敛', 'Debug 是否必须回流 Review', '人工 Gate 是否放在高风险处'], accent=RED, fill=RGBColor(0x1E,0x2A,0x4A), title_size=18, body_size=13, title_color=WHITE, body_color=RGBColor(0xD6,0xE4,0xF0))
    for idx,(t,c) in enumerate([('探索',TEAL),('规划',CYAN),('生成',GOLD),('审核',RED),('调试',GREEN)]):
        pill(slide, 10.0, 1.85 + idx*0.72, 1.7, 0.38, t, c, txtc=(NAVY if c in [CYAN,GOLD] else WHITE), size=14)

    # 2 boundaries
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
    textbox(slide, 0.7, 0.45, 6.2, 0.45, '系统边界与非目标', size=28, color=NAVY, bold=True)
    card(slide, 0.8, 1.2, 3.0, 2.15, '适用场景', ['RISC-V / KVM / 跨架构差距收敛', 'Kconfig / defconfig / selftest / kunit 一致性问题', '小步、可回滚、可解释的改动'], accent=TEAL)
    card(slide, 4.0, 1.2, 3.0, 2.15, '非目标', ['不追求无人监督自动发 patch', '不把“编过了”当作最终成功', '不允许单个 Agent 包办全流程'], accent=RED)
    card(slide, 7.2, 1.2, 2.45, 2.15, '成功标准', ['问题真实', '方案最小', '日志完整', 'patch-ready'], accent=GOLD)
    card(slide, 9.9, 1.2, 2.45, 2.15, '关键约束', ['每轮只跑 1 个 issue', '主实现单写分支', '审核独立运行'], accent=CYAN)
    codebox(slide, 0.85, 3.8, 5.8, 2.55, 'Design doctrine', '探索 -> 规划 -> 生成 -> 审核 -> 调试 -> patch-ready\n\n规则:\n1. 全流程工件化\n2. 生成/审核允许多轮\n3. Debug 不得直接放行\n4. 所有外部写操作前必须人工 Gate', accent=TEAL)
    card(slide, 7.0, 3.8, 5.3, 2.55, '为什么不是单模型串到底', ['内核上下文高度依赖历史与架构语义', '实现者自审会放大遗漏风险', '把失败拆给 Failure-Analyzer / Fix-Agent 更利于归因', '并行审核比“实现者复读”更稳'], accent=NAVY)

    # 3 state machine
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
    textbox(slide, 0.7, 0.45, 7.0, 0.45, '控制面状态机与停止条件', size=28, color=NAVY, bold=True)
    pill(slide, 0.95, 1.55, 1.7, 0.46, 'discover', TEAL)
    pill(slide, 3.0, 1.55, 1.5, 0.46, 'plan', CYAN, txtc=NAVY)
    pill(slide, 4.85, 1.55, 1.8, 0.46, 'generate', GOLD, txtc=NAVY)
    pill(slide, 7.0, 1.55, 1.6, 0.46, 'review', RED)
    pill(slide, 8.95, 1.55, 1.6, 0.46, 'debug', GREEN)
    pill(slide, 10.9, 1.55, 1.45, 0.46, 'patch', NAVY)
    for a,b,c in [(2.67,3.0,TEAL),(4.52,4.85,CYAN),(6.67,7.0,GOLD),(8.62,8.95,RED),(10.57,10.9,GREEN)]:
        arrow(slide, a, 1.78, b, 1.78, color=c)
    arrow(slide, 8.95, 2.2, 5.75, 2.95, color=GREEN)
    arrow(slide, 5.75, 2.95, 5.75, 2.02, color=GREEN)
    textbox(slide, 5.95, 2.62, 2.4, 0.28, 'debug -> review/generate 回流', size=12, color=GREEN)
    codebox(slide, 0.85, 3.4, 4.1, 2.6, 'workflow.yaml stop_conditions', 'max_rounds_generation_review: 3\nmax_rounds_debug: 2\nstop if decision == PATCH_READY\nstop if decision == BLOCKED\nstop if generation_review_round > 3\nstop if debug_round > 2', accent=RED)
    card(slide, 5.25, 3.4, 3.3, 2.6, '状态回写字段', ['current_stage', 'generation_review_round', 'debug_round', 'latest_decision', 'human_gate_status', 'artifact_status'], accent=CYAN)
    card(slide, 8.8, 3.4, 3.45, 2.6, '评审结论语义', ['PASS: 可进入下一阶段', 'REVISE / REVISE_GENERATION: 返回实现侧', 'ENTER_DEBUG: 进入失败归因', 'BLOCKED: 停机等人类判断'], accent=GOLD)

    # 4 role matrix
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
    textbox(slide, 0.7, 0.45, 6.8, 0.45, '角色矩阵：谁负责什么，不负责什么', size=28, color=NAVY, bold=True)
    card(slide, 0.8, 1.2, 2.35, 4.9, 'Orchestrator', ['负责状态机、工件路径、预算与调度', '不做架构结论', '不代替审核做 PASS 判定'], accent=NAVY)
    card(slide, 3.35, 1.2, 2.35, 4.9, 'Planner', ['产出 file-level 方案 / 测试矩阵 / 风险', '定义 must-not-do 边界', '不直接改代码'], accent=CYAN)
    card(slide, 5.9, 1.2, 2.35, 4.9, 'Implementer', ['按计划最小实现', '写构建/测试日志', '被审核驳回后进入再实现'], accent=GOLD)
    card(slide, 8.45, 1.2, 2.35, 4.9, 'Review trio', ['Spec: 看目标闭环', 'Code: 看质量/边界', 'Upstream: 看 patch 粒度与叙事'], accent=RED)
    card(slide, 11.0, 1.2, 1.55, 4.9, 'Debug', ['Failure-Analyzer', 'Fix-Agent', 'Regression-Guard'], accent=GREEN, title_size=17, body_size=12)
    textbox(slide, 0.95, 6.35, 11.2, 0.38, '技术判断：把“实现”和“判定通过”拆给不同 Agent，是这套系统避免自证正确的核心。', size=15, color=TEXT)

    # 5 artifact contract
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
    textbox(slide, 0.7, 0.45, 6.8, 0.45, '工件契约：让上下文留在文件，不留在聊天里', size=27, color=NAVY, bold=True)
    card(slide, 0.8, 1.2, 3.0, 2.25, 'Discover 工件', ['discover/code-evidence.md', 'discover/history-evidence.md', 'state/gap_registry.yaml', 'state/issue_map.yaml'], accent=TEAL)
    card(slide, 4.0, 1.2, 2.9, 2.25, 'Plan 工件', ['plans/design.md', 'plans/test-matrix.md', 'plans/risk.md'], accent=CYAN)
    card(slide, 7.1, 1.2, 2.8, 2.25, 'Generate 工件', ['logs/build-round-N.log', 'logs/test-round-N.log', 'state/change-summary-round-N.md'], accent=GOLD)
    card(slide, 10.1, 1.2, 2.35, 2.25, 'Patch 工件', ['series-cover-letter.md', 'checkpatch.txt', 'get-maintainer.txt'], accent=RED)
    codebox(slide, 0.9, 3.85, 5.2, 2.15, '目录约定', 'artifacts/<issue>/{discover, plans, review, debug, patch, state, logs}', accent=NAVY)
    codebox(slide, 6.35, 3.85, 5.95, 2.15, '命名约定', 'run_history/<issue>-round-<n>.md\nreview/spec-round-<n>.md\ndebug/failure-analysis-round-<n>.md\npatches/<issue>/draft-round-<n>.patch', accent=TEAL)

    # 6 review/debug loop
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
    textbox(slide, 0.7, 0.45, 6.8, 0.45, '生成 / 审核 / 调试 闭环设计', size=28, color=NAVY, bold=True)
    pill(slide, 1.0, 1.7, 2.0, 0.5, 'Generate R1', GOLD, txtc=NAVY, size=18)
    pill(slide, 3.6, 1.7, 2.0, 0.5, 'Spec Review', RED, size=18)
    pill(slide, 6.2, 1.7, 2.0, 0.5, 'Failure Analyze', GREEN, size=17)
    pill(slide, 8.8, 1.7, 1.85, 0.5, 'Fix Agent', TEAL, size=17)
    pill(slide, 11.1, 1.7, 1.2, 0.5, 'R2', NAVY, size=18)
    arrow(slide, 3.05, 1.95, 3.6, 1.95, color=GOLD)
    arrow(slide, 5.65, 1.95, 6.2, 1.95, color=RED)
    arrow(slide, 8.25, 1.95, 8.8, 1.95, color=GREEN)
    arrow(slide, 10.7, 1.95, 11.1, 1.95, color=TEAL)
    arrow(slide, 11.7, 2.22, 11.7, 3.15, color=CYAN)
    arrow(slide, 11.7, 3.15, 2.0, 3.15, color=CYAN)
    arrow(slide, 2.0, 3.15, 2.0, 2.22, color=CYAN)
    textbox(slide, 3.3, 2.92, 6.8, 0.3, '若审查仍未通过，则继续 R3；若长期不收敛，则触发人工 Gate。', size=13, color=CYAN, align=PP_ALIGN.CENTER)
    card(slide, 0.95, 4.15, 3.7, 1.85, '为什么要单独 Debug', ['很多失败不是编译错，而是规格未闭环', 'Failure-Analyzer 更像根因分析器，不是修代码器'], accent=GREEN)
    card(slide, 4.8, 4.15, 3.7, 1.85, 'Round 记录必须包含', ['本轮目标', '修改文件', '测试结果', 'review decision', 'next step'], accent=TEAL)
    card(slide, 8.65, 4.15, 3.7, 1.85, '停止条件', ['连续 3 轮不收敛', '错误反复出现且根因不稳定', '出现架构语义争议'], accent=RED)

    # 7 quality gates
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
    textbox(slide, 0.7, 0.45, 7.2, 0.45, '质量闸门与验证矩阵', size=28, color=NAVY, bold=True)
    stat_card(slide, 0.85, 1.15, 2.2, 1.35, '3', '人工 Gate', RED)
    stat_card(slide, 3.25, 1.15, 2.2, 1.35, '2+', '生成/审核轮次', CYAN)
    stat_card(slide, 5.65, 1.15, 2.2, 1.35, '1', 'issue per run', GOLD)
    stat_card(slide, 8.05, 1.15, 2.2, 1.35, '6+', '关键工件', TEAL)
    card(slide, 0.9, 3.0, 3.8, 2.8, '自动质量闸门', ['make ARCH=riscv defconfig', '目标架构 build', 'kselftest relevant subset', 'kunit relevant subset', 'checkpatch.pl', 'get_maintainer.pl'], accent=GOLD)
    card(slide, 4.95, 3.0, 3.8, 2.8, '人工 Gate 触发条件', ['问题定义不清', 'ABI/UAPI/DT/Kconfig 用户可见影响', 'maintainer rejection history', '准备发信'], accent=RED)
    card(slide, 9.0, 3.0, 3.35, 2.8, '评审视角', ['Spec: 需求是否闭环', 'Code: 边界与可维护性', 'Upstream: patch 叙事/粒度'], accent=CYAN)

    # 8 risks/costs
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
    textbox(slide, 0.7, 0.45, 6.8, 0.45, '成本、风险与工程控制', size=28, color=NAVY, bold=True)
    card(slide, 0.8, 1.2, 3.0, 4.8, '典型失败模式', ['把设计问题误判成代码问题', '实现者过度保守只改一半', '过度设计导致 patch 粒度膨胀', '长日志回灌导致 token 成本失控'], accent=RED)
    card(slide, 4.05, 1.2, 3.0, 4.8, '成本控制策略', ['每次只处理一个窄 issue', '长上下文外置为工件路径', '高歧义规划交给更强模型', '实现闭环尽量在 cheaper runtime 完成'], accent=TEAL)
    card(slide, 7.3, 1.2, 2.4, 4.8, '收敛策略', ['round 上限', 'debug 上限', '人工 Gate', '单 issue 单分支'], accent=CYAN)
    card(slide, 9.95, 1.2, 2.35, 4.8, '最终价值', ['把人工从重复劳动里解放出来', '保留关键判断', '形成可复盘工程资产'], accent=GOLD)

    # 9 demo
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
    textbox(slide, 0.7, 0.45, 8.2, 0.45, 'Demo Issue：riscv/kvm MMU Kconfig 一致性修复', size=27, color=NAVY, bold=True)
    card(slide, 0.8, 1.15, 2.7, 4.8, 'Issue', ['目标：让 dependency/help text 都显式表达 MMU 约束', '范围：只触及 arch/riscv/kvm/Kconfig', '风险：低，但属于 Kconfig 用户可见变化'], accent=NAVY)
    card(slide, 3.7, 1.15, 2.7, 4.8, 'Round 1', ['实现通过模拟构建检查', '但 Spec Review 指出 help text 未闭环', '结论：ENTER_DEBUG'], accent=GOLD)
    card(slide, 6.6, 1.15, 2.7, 4.8, 'Debug', ['Failure-Analyzer 认定为 spec gap', 'Fix-Agent 继续单文件修复', '保持单 patch、可回滚'], accent=RED)
    card(slide, 9.5, 1.15, 2.8, 4.8, 'Round 2', ['补充 help text 的 MMU 说明', 'Spec Review Round 2: PASS', '状态流转到 patch-ready'], accent=GREEN)
    textbox(slide, 0.95, 6.35, 11.3, 0.38, '这个 demo 证明：审查失败并不一定意味着代码坏了，也可能只是“实现没有完整满足设计目标”。', size=14, color=TEXT)

    # 10 conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, NAVY)
    textbox(slide, 0.8, 0.7, 6.8, 0.5, '技术评审结论', size=30, color=WHITE, bold=True)
    card(slide, 0.95, 1.55, 3.55, 3.8, '这套方案为什么可行', ['职责解耦：探索、规划、实现、审核、调试分层', '状态显式：所有阶段都回写 workflow.yaml 与工件目录', '闭环可控：round 上限 + 人工 Gate + patch-ready 定义'], accent=TEAL, fill=RGBColor(0xF5,0xF8,0xFC))
    card(slide, 4.9, 1.55, 3.55, 3.8, '仍需谨慎的地方', ['架构语义争议', '真实 lore 证据补齐', '日志真实性与回归覆盖', '发信礼仪与 maintainer 偏好'], accent=RED, fill=RGBColor(0xF5,0xF8,0xFC))
    card(slide, 8.85, 1.55, 3.45, 3.8, '建议决策', ['先用 MVP 跑一个高置信度窄 issue', '强制经历至少 2 轮生成/审核', '把真实工具链日志接入现有工件目录'], accent=GOLD, fill=RGBColor(0xF5,0xF8,0xFC))
    pill(slide, 3.95, 6.0, 5.4, 0.58, '不是“多开模型”，而是“多角色互相制衡”', GOLD, txtc=NAVY, size=19)

    out = Path('ai-assisted-kernel-development-tech-review.pptx')
    prs.save(out)
    return out


def make_customized():
    prs = new_prs('AI辅助内核开发多智能体工作流：定制版', 'customized with local docs and artifacts')
    workflow_excerpt = snippet('workflow.kernel-multi-agent.example.yaml', 0, 18)
    issue_excerpt = '\n'.join([
        '# Issue Brief',
        '',
        'Issue ID: demo-riscv-kvm-mmu-kconfig',
        'Title: riscv/kvm: explicit MMU requirement',
        'Type: low-risk Kconfig consistency fix',
        'Subsystem: RISC-V / KVM',
        'Desired outcome: dependency/help text 与 host KVM 语义一致',
    ])
    design_excerpt = '\n'.join([
        'Minimal change path',
        '1. 只修改 arch/riscv/kvm/Kconfig',
        '2. dependency + help text 都显式体现 MMU 约束',
        '3. 不修改运行时代码',
        '4. 不引入新符号 / 不做无关 cleanup',
        'Review only: arch/riscv/Kconfig, defconfig',
    ])
    cover_excerpt = '\n'.join([
        'Subject: [PATCH 0/1] riscv/kvm: clarify MMU requirement',
        '',
        'Single low-risk Kconfig consistency fix for RISC-V host KVM.',
        'No runtime behavior change.',
        'Makes MMU requirement explicit in dependency/help text.',
    ])
    fail_excerpt = snippet('artifacts/demo-riscv-kvm-mmu-kconfig/debug/failure-analysis-round-1.md', 8, 24)

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, NAVY)
    textbox(slide, 0.8, 1.0, 8.4, 1.2, 'AI辅助内核开发\n多智能体工作流（定制版）', size=28, color=WHITE, bold=True)
    textbox(slide, 0.82, 2.4, 8.9, 1.0, '基于当前工作区真实工件定制：workflow 文档、workflow.yaml、demo issue、debug 记录、cover letter 与 checkpatch 输出。', size=17, color=RGBColor(0xD6,0xE4,0xF0))
    pill(slide, 0.82, 3.58, 2.2, 0.44, 'Local Docs', TEAL)
    pill(slide, 3.2, 3.58, 2.2, 0.44, 'Real Artifact Paths', CYAN, txtc=NAVY)
    pill(slide, 5.58, 3.58, 2.4, 0.44, 'Demo Patch-Ready', GOLD, txtc=NAVY)
    card(slide, 9.45, 1.1, 2.8, 4.9, '本次引用文件', ['ai-assisted-kernel-development-multi-agent-workflow.md', 'workflow.kernel-multi-agent.example.yaml', 'artifacts/demo-riscv-kvm-mmu-kconfig/...'], accent=TEAL, fill=RGBColor(0x1E,0x2A,0x4A), title_size=18, body_size=12, title_color=WHITE, body_color=RGBColor(0xD6,0xE4,0xF0))

    # 2 source map
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
    textbox(slide, 0.7, 0.45, 6.5, 0.45, '定制输入材料与落盘位置', size=28, color=NAVY, bold=True)
    card(slide, 0.8, 1.2, 3.0, 2.1, '方案文档', ['workflow 主文档', 'gap workflow 文档', 'OpenClaw Linux RISC-V 方案'], accent=TEAL)
    card(slide, 4.0, 1.2, 3.0, 2.1, '控制面配置', ['workflow.kernel-multi-agent.example.yaml', 'artifacts/.../state/workflow.yaml'], accent=CYAN)
    card(slide, 7.2, 1.2, 2.5, 2.1, '问题工件', ['issue-brief.md', 'plans/design.md', 'review/spec-round-2.md'], accent=GOLD)
    card(slide, 9.95, 1.2, 2.35, 2.1, 'Patch 工件', ['series-cover-letter.md', 'checkpatch.txt', 'get-maintainer.txt'], accent=RED)
    codebox(slide, 0.9, 3.8, 5.2, 2.2, 'workspace root', 'artifacts/demo-riscv-kvm-mmu-kconfig/\n  discover/ plans/ review/\n  debug/ patch/ state/ logs/', accent=NAVY)
    card(slide, 6.45, 3.8, 5.8, 2.2, '为什么这版更“定制”', ['不再只讲抽象架构，而是直接映射到当前仓库中的真实路径、真实 issue_id 和真实工件名称', '汇报时可直接点击到对应文件继续深挖'], accent=TEAL)

    # 3 local architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
    textbox(slide, 0.7, 0.45, 7.0, 0.45, '本地架构图：OpenClaw / Claude Code / Codex / Human', size=27, color=NAVY, bold=True)
    pill(slide, 5.55, 0.95, 2.2, 0.44, 'Human Gates', RED)
    pill(slide, 5.35, 1.65, 2.6, 0.54, 'OpenClaw 控制面', NAVY, size=20)
    card(slide, 0.9, 2.5, 2.9, 2.65, '输入源', ['torvalds/linux 工作树', 'KVM lore / 历史讨论', 'linux-riscv-docs issues'], accent=TEAL)
    card(slide, 4.35, 2.5, 2.9, 2.65, 'Claude Code 侧', ['Scout / Planner / Review / Failure-Analyzer', '擅长：方案、审查、归因'], accent=CYAN)
    card(slide, 8.0, 2.5, 2.9, 2.65, 'Codex 侧', ['Implementer / Fix-Agent', '擅长：实现、构建、测试、修补'], accent=GOLD)
    card(slide, 11.15, 2.5, 1.2, 2.65, '产出', ['plans', 'logs', 'patch', 'review'], accent=GREEN, title_size=17, body_size=12)
    arrow(slide, 3.8, 3.1, 5.35, 1.92, color=TEAL)
    arrow(slide, 6.65, 2.2, 5.9, 2.5, color=CYAN)
    arrow(slide, 7.95, 1.92, 8.0, 3.1, color=GOLD)
    arrow(slide, 10.9, 3.8, 11.15, 3.8, color=GREEN)
    textbox(slide, 0.95, 6.1, 11.2, 0.42, '这张图直接对应本仓库文档中的“控制面 + 执行 Agent + 人工 Gate”设计，而不是泛化 AI 工作流模板。', size=14, color=TEXT)

    # 4 workflow yaml excerpt
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
    textbox(slide, 0.7, 0.45, 6.7, 0.45, '现有 workflow.yaml 片段：状态机真正规约在哪里', size=27, color=NAVY, bold=True)
    codebox(slide, 0.85, 1.15, 6.0, 4.95, 'workflow.yaml excerpt', workflow_excerpt, accent=TEAL)
    card(slide, 7.1, 1.2, 5.1, 1.55, '读法', ['workspace 决定工件目录；controller 决定重试与 stop；human_gates 决定何时必须人工审批'], accent=CYAN)
    card(slide, 7.1, 2.95, 5.1, 1.45, '这意味着什么', ['Agent 可以被替换，但工件契约与状态机不应漂移', '系统稳定性来自文件规范，不来自单次 prompt 灵感'], accent=GOLD)
    card(slide, 7.1, 4.65, 5.1, 1.45, '汇报价值', ['领导可看 Gate 与风险；工程师可看 outputs 与 stop_conditions；执行者可看 agent inputs/outputs'], accent=RED)

    # 5 issue + design
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
    textbox(slide, 0.7, 0.45, 8.0, 0.45, '真实示例：Issue Brief + Design Plan', size=27, color=NAVY, bold=True)
    codebox(slide, 0.85, 1.15, 5.7, 2.25, 'issue-brief.md', issue_excerpt, accent=NAVY, size=10)
    codebox(slide, 0.85, 3.75, 5.7, 2.2, 'design.md (minimal path)', design_excerpt, accent=CYAN, size=10)
    card(slide, 6.9, 1.2, 5.2, 4.75, '这页想传达的关键信息', ['这不是“随便找个 bug”，而是一个被明确定义为低风险、单文件、可回滚的 demo issue', 'Design plan 明确写出 must-not-do boundary：不引入新符号、不碰 runtime、不做无关 cleanup', '这类 issue 最适合做 MVP 首轮演示'], accent=TEAL)

    # 6 debug evidence
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, LIGHT)
    textbox(slide, 0.7, 0.45, 7.2, 0.45, '真实示例：为什么 Round 1 被打回', size=27, color=NAVY, bold=True)
    codebox(slide, 0.85, 1.15, 6.2, 4.85, 'failure-analysis-round-1.md', fail_excerpt, accent=RED)
    card(slide, 7.35, 1.2, 4.9, 1.75, '根因不是编译错误', ['构建/配置检查没炸，但规格目标要求“dependency + help text 一并澄清”', 'Round 1 只完成了一半'], accent=GOLD)
    card(slide, 7.35, 3.15, 4.9, 1.35, '因此需要 Debug', ['Debug 在这里承担“规格缺口归因”而不是“修一个 crash”'], accent=GREEN)
    card(slide, 7.35, 4.75, 4.9, 1.25, '汇报亮点', ['非常适合向听众解释：为什么生成和审核一定要允许多轮迭代'], accent=TEAL)

    # 7 patch ready artifacts
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, WHITE)
    textbox(slide, 0.7, 0.45, 8.1, 0.45, 'Patch-ready 不是一句话，而是一组真实工件', size=27, color=NAVY, bold=True)
    codebox(slide, 0.85, 1.15, 5.8, 2.65, 'series-cover-letter.md', cover_excerpt, accent=NAVY)
    card(slide, 6.95, 1.2, 5.15, 1.65, 'checkpatch 结果', ['artifacts/.../patch/checkpatch.txt', '当前 demo: total 0 errors, 0 warnings, 0 checks'], accent=GREEN)
    card(slide, 6.95, 3.1, 5.15, 1.35, 'get_maintainer 输出', ['RISC-V + KVM maintainers', 'linux-riscv@lists.infradead.org', 'kvm@vger.kernel.org'], accent=CYAN)
    card(slide, 6.95, 4.7, 5.15, 1.3, '最后仍需 Gate-3', ['这版定制汇报会明确说明：patch-ready ≠ 自动发送'], accent=RED)

    # 8 implementation recommendation
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, NAVY)
    textbox(slide, 0.8, 0.65, 7.0, 0.5, '怎么把这版继续用在你的实际汇报里', size=29, color=WHITE, bold=True)
    card(slide, 0.9, 1.55, 3.6, 3.8, '适合技术评审会', ['展示 workflow.yaml 片段', '展示 issue/design/debug/patch 实际路径', '强调多轮审查与人工 Gate'], accent=TEAL, fill=RGBColor(0xF5,0xF8,0xFC))
    card(slide, 4.85, 1.55, 3.6, 3.8, '适合项目推进会', ['把 demo issue 换成你的下一条真实 issue', '把 simulated build/test 替换成真实日志', '补进真实 get_maintainer / lore thread'], accent=CYAN, fill=RGBColor(0xF5,0xF8,0xFC))
    card(slide, 8.8, 1.55, 3.45, 3.8, '下一步建议', ['增加真实 patch diff 截图或代码片段', '加入成本统计页（轮次 / token / 时长）', '把架构图导出成正式 SVG 后替换'], accent=GOLD, fill=RGBColor(0xF5,0xF8,0xFC))
    pill(slide, 4.0, 6.0, 5.1, 0.58, '这版已经绑定到当前工作区的真实文件结构', GOLD, txtc=NAVY, size=19)

    out = Path('ai-assisted-kernel-development-customized.pptx')
    prs.save(out)
    return out


if __name__ == '__main__':
    a = make_tech_review()
    b = make_customized()
    print(a)
    print(b)
