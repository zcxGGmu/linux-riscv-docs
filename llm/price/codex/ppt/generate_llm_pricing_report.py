from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(18, 49, 91)
BLUE = RGBColor(51, 102, 204)
LIGHT_BLUE = RGBColor(230, 240, 252)
ORANGE = RGBColor(233, 131, 0)
LIGHT_ORANGE = RGBColor(253, 241, 223)
TEAL = RGBColor(47, 132, 125)
LIGHT_TEAL = RGBColor(228, 244, 242)
GRAY = RGBColor(92, 101, 112)
LIGHT_GRAY = RGBColor(243, 246, 249)
DARK = RGBColor(31, 41, 55)
WHITE = RGBColor(255, 255, 255)

FONT_FAMILY = "Microsoft YaHei"


def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_title_slide(prs)
    _add_exec_summary_slide(prs)
    _add_pricing_slide(prs)
    _add_limits_slide(prs)
    _add_token_budget_slide(prs)
    _add_usage_breakdown_slide(prs)
    _add_cost_slide(prs)
    _add_routing_slide(prs)
    _add_cost_control_slide(prs)
    _add_rollout_slide(prs)
    _add_risk_slide(prs)
    _add_appendix_slide(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def _add_header(slide, title: str, subtitle: str | None = None) -> None:
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.55))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY

    title_shape = slide.shapes.title
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.72)
    title_shape.width = Inches(8.4)
    title_shape.height = Inches(0.6)
    p = title_shape.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT_FAMILY
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.22), Inches(10.8), Inches(0.45))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT_FAMILY
        sp.font.size = Pt(10)
        sp.font.color.rgb = GRAY


def _add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.03), Inches(12.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT_FAMILY
    p.font.size = Pt(8)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT


def _add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.bullet = True
        p.font.name = FONT_FAMILY
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
        p.line_spacing = 1.15


def _add_card(slide, left: float, top: float, width: float, height: float, title: str, value: str, note: str, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.14), width - Inches(0.36), Inches(0.28))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = title
    tp.font.name = FONT_FAMILY
    tp.font.size = Pt(12)
    tp.font.bold = True
    tp.font.color.rgb = NAVY

    value_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.48), width - Inches(0.36), Inches(0.5))
    vp = value_box.text_frame.paragraphs[0]
    vp.text = value
    vp.font.name = FONT_FAMILY
    vp.font.size = Pt(24)
    vp.font.bold = True
    vp.font.color.rgb = DARK

    note_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(1.02), width - Inches(0.36), Inches(0.52))
    np = note_box.text_frame.paragraphs[0]
    np.text = note
    np.font.name = FONT_FAMILY
    np.font.size = Pt(10)
    np.font.color.rgb = GRAY


def _style_cell(cell, *, bold: bool = False, fill: RGBColor | None = None, font_size: int = 11, align=PP_ALIGN.CENTER) -> None:
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = DARK if fill != NAVY else WHITE
    if not cell.text_frame.paragraphs[0].runs:
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = cell.text
        run.font.name = FONT_FAMILY
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = DARK if fill != NAVY else WHITE


def _add_title_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.color.rgb = NAVY

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(1.28), Inches(1.7), Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ORANGE
    accent.line.color.rgb = ORANGE

    title = slide.shapes.add_textbox(Inches(0.62), Inches(1.48), Inches(11.4), Inches(1.55))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "大模型定价与 20 人团队\n月度 Token 用量评估"
    p.font.name = FONT_FAMILY
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    subtitle = slide.shapes.add_textbox(Inches(0.62), Inches(3.32), Inches(8.6), Inches(1.0))
    sp = subtitle.text_frame.paragraphs[0]
    sp.text = "覆盖 GLM-5.1 / Claude Opus 4.6 / GPT-5.4\n输出正式汇报口径、预算区间与落地建议"
    sp.font.name = FONT_FAMILY
    sp.font.size = Pt(16)
    sp.font.color.rgb = RGBColor(220, 228, 240)

    meta = slide.shapes.add_textbox(Inches(0.62), Inches(6.15), Inches(4.2), Inches(0.5))
    mp = meta.text_frame.paragraphs[0]
    mp.text = "版本：2026-04-15"
    mp.font.name = FONT_FAMILY
    mp.font.size = Pt(12)
    mp.font.color.rgb = RGBColor(220, 228, 240)


def _add_exec_summary_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "执行摘要", "面向管理层的 4 个决定性结论")

    _add_card(slide, Inches(0.5), Inches(1.8), Inches(3.0), Inches(1.7), "预算基线", "6.18 亿", "20 人团队月度 token 基线预算", LIGHT_BLUE)
    _add_card(slide, Inches(3.7), Inches(1.8), Inches(3.0), Inches(1.7), "主成本池", "57%", "代码开发 + Review loop 占总量", LIGHT_ORANGE)
    _add_card(slide, Inches(6.9), Inches(1.8), Inches(3.0), Inches(1.7), "最低人民币口径", "GLM-5.1", "但长上下文会触发更高分档", LIGHT_TEAL)
    _add_card(slide, Inches(10.1), Inches(1.8), Inches(2.7), Inches(1.7), "国际主力建议", "GPT-5.4", "Opus 4.6 适合高价值升级任务", LIGHT_GRAY)

    bullets = [
        "GLM-5.1 当前公开人民币单价最低，但输入长度 <32K / >=32K 分档会直接影响成本。",
        "GPT-5.4 的价格、缓存输入、超长上下文能力组合最均衡，适合作为国际通用主力模型。",
        "Claude Opus 4.6 适合复杂架构推演、关键疑难问题定位和最终高价值 Review，不适合全员全程默认主力。",
        "更合理的规划不是只比单价，而是按“主力模型 + 升级模型 + 成本控制策略”组合决策。",
    ]
    _add_bullets(slide, Inches(0.68), Inches(3.95), Inches(11.8), Inches(2.6), bullets)
    _add_footer(slide, "来源：OpenAI / Anthropic / 智谱官方定价与模型文档，2026-04-15")


def _add_pricing_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "三家模型价格总览", "先看公开 API 单价，再看缓存和长上下文门槛")

    summary = slide.shapes.add_textbox(Inches(0.55), Inches(1.6), Inches(12.1), Inches(0.5))
    sp = summary.text_frame.paragraphs[0]
    sp.text = "结论：GLM-5.1 人民币口径最低；GPT-5.4 国际价格最均衡；Claude Opus 4.6 适合高价值升级而非全量承接。"
    sp.font.name = FONT_FAMILY
    sp.font.size = Pt(14)
    sp.font.bold = True
    sp.font.color.rgb = NAVY

    rows, cols = 4, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(2.15), Inches(12.15), Inches(2.8)).table
    headers = ["模型", "输入价格", "缓存/命中", "输出价格", "备注"]
    data = [
        ["GLM-5.1", "6 / 8 元\n(<32K / >=32K)", "1.3 / 2 元\n命中价", "24 / 28 元", "缓存存储限时免费\n有效期至 8 月 31 日"],
        ["Claude Opus 4.6", "$5 / MTok", "缓存命中最高节省 90%", "$25 / MTok", "Batch 50% 折扣\n1M beta 上下文更贵"],
        ["GPT-5.4", "$2.50 / 1M", "$0.25 / 1M\ncached input", "$15.00 / 1M", ">272K 后输入 2x\n输出 1.5x"],
    ]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
        _style_cell(table.cell(0, idx), bold=True, fill=NAVY, font_size=11)
    for r, row in enumerate(data, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
            _style_cell(table.cell(r, c), fill=WHITE, font_size=11)

    _add_card(slide, Inches(0.7), Inches(5.25), Inches(3.7), Inches(1.2), "最低人民币成本", "GLM-5.1", "适合中文研发协作与成本敏感场景", LIGHT_BLUE)
    _add_card(slide, Inches(4.8), Inches(5.25), Inches(3.7), Inches(1.2), "最均衡国际主力", "GPT-5.4", "缓存输入清晰、上下文能力强", LIGHT_TEAL)
    _add_card(slide, Inches(8.9), Inches(5.25), Inches(3.3), Inches(1.2), "高质量升级", "Opus 4.6", "适合复杂任务，不适合全量默认", LIGHT_ORANGE)
    _add_footer(slide, "价格口径：公开 API 定价；未计入税费、企业协议折扣和代理层加价")


def _add_limits_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "三家模型限制总览", "真正影响可用性的不是单价，而是上下文、最大输出和并发")

    rows, cols = 4, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.65), Inches(12.15), Inches(3.1)).table
    headers = ["模型", "上下文窗口", "最大输出", "速率/并发", "备注"]
    data = [
        ["GLM-5.1", "200K", "128K", "官方未公开固定 RPM/TPM\n按权益与并发管理", "Coding Plan 并发不可单独申请调整"],
        ["Claude Opus 4.6", "200K 标准\n1M beta", "128K", "Tier 1: 50 RPM\n30,000 ITPM / 8,000 OTPM", "高价值复杂任务更合适"],
        ["GPT-5.4", "1,048,576", "128K", "Tier 1: 500 RPM / 500K TPM\nTier 4: 10K RPM / 4M TPM", "公开规格最完整也最透明"],
    ]
    for idx, header in enumerate(headers):
        table.cell(0, idx).text = header
        _style_cell(table.cell(0, idx), bold=True, fill=NAVY, font_size=11)
    for r, row in enumerate(data, start=1):
        for c, value in enumerate(row):
            table.cell(r, c).text = value
            _style_cell(table.cell(r, c), fill=WHITE, font_size=11)

    bullets = [
        "长上下文能力最强：GPT-5.4。",
        "复杂推理升级通道：Claude Opus 4.6。",
        "团队规模化成本优势：GLM-5.1，但要控制 >=32K 输入档占比。",
    ]
    _add_bullets(slide, Inches(0.7), Inches(5.1), Inches(6.2), Inches(1.4), bullets)

    tip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(5.0), Inches(5.1), Inches(1.45))
    tip.fill.solid()
    tip.fill.fore_color.rgb = LIGHT_ORANGE
    tip.line.color.rgb = LIGHT_ORANGE
    note = slide.shapes.add_textbox(Inches(7.52), Inches(5.18), Inches(4.7), Inches(1.0))
    np = note.text_frame.paragraphs[0]
    np.text = "注意：GLM-5.1 模型页示例代码出现 max_tokens=65536，但官方模型概览页写明最大输出为 128K，本汇报以概览页为准。"
    np.font.name = FONT_FAMILY
    np.font.size = Pt(10)
    np.font.color.rgb = DARK
    _add_footer(slide, "限制口径：官方模型页 / 模型概览 / 速率限制文档")


def _add_token_budget_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "20 人团队月度 Token 基线模型", "预算区间的关键不是点值，而是基线与波动范围")

    chart_data = CategoryChartData()
    chart_data.categories = ["保守", "基线", "激进"]
    chart_data.add_series("总量(百万)", (370.9, 618.2, 865.5))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.7),
        Inches(1.8),
        Inches(6.4),
        Inches(3.8),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLUE

    _add_card(slide, Inches(7.6), Inches(1.9), Inches(2.2), Inches(1.55), "保守", "370.9M", "AI 辅助问答 + 有限 Review", LIGHT_GRAY)
    _add_card(slide, Inches(9.95), Inches(1.9), Inches(2.2), Inches(1.55), "基线", "618.2M", "AI 深度接入开发与测试循环", LIGHT_BLUE)
    _add_card(slide, Inches(7.6), Inches(3.65), Inches(4.55), Inches(1.55), "激进", "865.5M", "高频 agentic coding + 长上下文 + 重度 Review", LIGHT_ORANGE)

    bullets = [
        "规划区间建议按 3.7 亿到 8.7 亿 tokens/月，而不是只预留一个静态点值。",
        "基线场景可作为预算、限流和路由策略的默认经营口径。",
        "如果团队开始大规模使用自动修复、长上下文代码代理，消耗会向激进场景快速靠拢。",
    ]
    _add_bullets(slide, Inches(0.72), Inches(5.95), Inches(11.6), Inches(0.9), bullets)
    _add_footer(slide, "假设：20 人、22 个工作日、需求/开发/测试三环节都深度使用大模型")


def _add_usage_breakdown_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "分环节用量拆分", "代码开发与 Review loop 是绝对主成本池")

    chart_data = CategoryChartData()
    chart_data.categories = ["需求生成", "代码开发+Review", "测试验证"]
    chart_data.add_series("新输入", (35.2, 79.2, 35.2))
    chart_data.add_series("缓存命中", (61.6, 184.8, 63.8))
    chart_data.add_series("输出", (26.4, 88.0, 44.0))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED,
        Inches(0.7),
        Inches(1.8),
        Inches(7.0),
        Inches(3.9),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(10)
    for series, color in zip(chart.series, (BLUE, ORANGE, TEAL)):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    _add_card(slide, Inches(8.25), Inches(1.9), Inches(4.0), Inches(1.35), "最大头", "56.9%", "代码开发 + Review 占比最高", LIGHT_ORANGE)
    _add_card(slide, Inches(8.25), Inches(3.45), Inches(4.0), Inches(1.35), "管理重点", "缓存复用", "高频往返和重复读取决定账单", LIGHT_TEAL)
    _add_card(slide, Inches(8.25), Inches(5.0), Inches(4.0), Inches(1.35), "不要误判", "不是需求阶段", "需求分析不是主成本来源", LIGHT_BLUE)

    _add_footer(slide, "基线拆分：需求 123.2M / 开发+Review 352.0M / 测试 143.0M")


def _add_cost_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "基线月费测算", "同一基线流量下，模型路由策略比单价对比更重要")

    _add_card(slide, Inches(0.75), Inches(1.9), Inches(3.75), Inches(2.25), "GLM-5.1", "5,102 元 ~ 6,252 元/月", "取决于 <32K / >=32K 输入长度分档；混合负载下更接近 5.6K 元/月", LIGHT_BLUE)
    _add_card(slide, Inches(4.78), Inches(1.9), Inches(3.75), Inches(2.25), "GPT-5.4", "$2,828 /月", "已计入 cached input 价格；若频繁超过 272K 上下文，成本会继续上浮", LIGHT_TEAL)
    _add_card(slide, Inches(8.8), Inches(1.9), Inches(3.75), Inches(2.25), "Claude Opus 4.6", "$4,863 ~ $6,259 /月", "适合作为升级模型；全量默认使用会显著放大账单", LIGHT_ORANGE)

    bullets = [
        "如果目标是人民币成本最低，优先考虑 GLM-5.1，但要严控长输入分档。",
        "如果目标是国际主力模型，GPT-5.4 的成本/能力组合最均衡。",
        "如果目标是难题突破率，保留 Claude Opus 4.6 作为升级通道，而不是日常默认流量入口。",
    ]
    _add_bullets(slide, Inches(0.72), Inches(4.55), Inches(11.8), Inches(1.55), bullets)
    _add_footer(slide, "说明：未计税费、企业协议折扣、代理层加价和私有部署费用")


def _add_routing_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "模型分工建议", "用“主力模型 + 升级模型”替代“单一最好模型”")

    boxes = [
        (Inches(0.8), LIGHT_BLUE, "日常开发主力", "GLM-5.1 / GPT-5.4", "普通编码、解释、样板测试、常规 Review"),
        (Inches(4.6), LIGHT_TEAL, "复杂任务升级", "Claude Opus 4.6", "架构争议、关键性能瓶颈、最终高价值 Review"),
        (Inches(8.4), LIGHT_ORANGE, "成本优化层", "缓存 + 裁剪 + 路由", "控制 >=32K 输入档和超长上下文触发频率"),
    ]
    for left, fill, title, model, note in boxes:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(3.2), Inches(2.35))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill

        tb = slide.shapes.add_textbox(left + Inches(0.18), Inches(2.38), Inches(2.8), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT_FAMILY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = NAVY

        mb = slide.shapes.add_textbox(left + Inches(0.18), Inches(2.86), Inches(2.8), Inches(0.45))
        mp = mb.text_frame.paragraphs[0]
        mp.text = model
        mp.font.name = FONT_FAMILY
        mp.font.size = Pt(18)
        mp.font.bold = True
        mp.font.color.rgb = DARK

        nb = slide.shapes.add_textbox(left + Inches(0.18), Inches(3.42), Inches(2.8), Inches(0.85))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.name = FONT_FAMILY
        np.font.size = Pt(11)
        np.font.color.rgb = GRAY

    for left in (Inches(3.85), Inches(7.65)):
        arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, Inches(2.9), Inches(0.45), Inches(0.55))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = BLUE
        arrow.line.color.rgb = BLUE

    bullets = [
        "国内优先：GLM-5.1 适合中文研发协作与预算敏感场景。",
        "国际优先：GPT-5.4 适合作为默认主力模型。",
        "不建议：把 Claude Opus 4.6 承接所有普通补全、样板代码和测试脚手架。",
    ]
    _add_bullets(slide, Inches(0.72), Inches(5.2), Inches(11.8), Inches(1.1), bullets)
    _add_footer(slide, "推荐策略：主力模型承接 80% 流量，升级模型处理 20% 的高价值难题")


def _add_cost_control_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "成本控制抓手", "先抓主成本池，再抓长上下文，再抓缓存复用")

    items = [
        ("抓住主成本池", "优先优化代码开发 + Review loop，而不是只盯需求分析。", LIGHT_BLUE),
        ("控制长上下文", "避免 GLM-5.1 大量落到 >=32K 输入档，避免 GPT-5.4 高频触发 >272K 计费门槛。", LIGHT_ORANGE),
        ("强制上下文裁剪", "只传相关文件片段、关键日志片段和最小可复现信息。", LIGHT_TEAL),
        ("开启缓存复用", "在 IDE / Agent 层优先复用上下文，减少重复读仓库和日志。", LIGHT_GRAY),
        ("做模型分级", "主力模型负责规模化吞吐，贵模型只承担升级任务。", LIGHT_BLUE),
    ]
    top = Inches(1.75)
    for idx, (title, note, fill) in enumerate(items):
        box_top = top + Inches(idx * 1.0)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.72), box_top, Inches(11.9), Inches(0.78))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill

        tbox = slide.shapes.add_textbox(Inches(0.92), box_top + Inches(0.1), Inches(2.2), Inches(0.25))
        tp = tbox.text_frame.paragraphs[0]
        tp.text = title
        tp.font.name = FONT_FAMILY
        tp.font.size = Pt(12)
        tp.font.bold = True
        tp.font.color.rgb = NAVY

        nbox = slide.shapes.add_textbox(Inches(3.05), box_top + Inches(0.1), Inches(9.1), Inches(0.45))
        np = nbox.text_frame.paragraphs[0]
        np.text = note
        np.font.name = FONT_FAMILY
        np.font.size = Pt(11)
        np.font.color.rgb = DARK
    _add_footer(slide, "控制顺序建议：开发 Review → 长上下文 → 缓存复用 → 路由分级")


def _add_rollout_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "推荐落地路径", "先试点，再放量，再按成本和质量做路由收敛")

    steps = [
        ("01", "先试点", "选 1-2 个研发小组，跑通需求分析、开发、测试三环节。"),
        ("02", "定口径", "统一记录新输入、缓存命中、输出和真实月账单。"),
        ("03", "做路由", "把主力模型与升级模型职责固定下来，减少人工切换。"),
        ("04", "看经营指标", "每月复盘成本、响应质量、交付效率和失败重试率。"),
    ]
    for idx, (num, title, note) in enumerate(steps):
        left = Inches(0.82 + idx * 3.05)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, Inches(2.2), Inches(0.72), Inches(0.72))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.color.rgb = BLUE
        cb = slide.shapes.add_textbox(left, Inches(2.33), Inches(0.72), Inches(0.2))
        cp = cb.text_frame.paragraphs[0]
        cp.text = num
        cp.font.name = FONT_FAMILY
        cp.font.size = Pt(12)
        cp.font.bold = True
        cp.font.color.rgb = WHITE
        cp.alignment = PP_ALIGN.CENTER

        tb = slide.shapes.add_textbox(left + Inches(0.92), Inches(2.16), Inches(1.75), Inches(0.3))
        tp = tb.text_frame.paragraphs[0]
        tp.text = title
        tp.font.name = FONT_FAMILY
        tp.font.size = Pt(15)
        tp.font.bold = True
        tp.font.color.rgb = NAVY

        nb = slide.shapes.add_textbox(left + Inches(0.92), Inches(2.55), Inches(1.95), Inches(1.0))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.name = FONT_FAMILY
        np.font.size = Pt(11)
        np.font.color.rgb = DARK
        np.word_wrap = True

        if idx < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left + Inches(2.62), Inches(2.42), Inches(0.38), Inches(0.38))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ORANGE
            arrow.line.color.rgb = ORANGE

    note_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(4.65), Inches(11.4), Inches(1.15))
    note_shape.fill.solid()
    note_shape.fill.fore_color.rgb = LIGHT_GRAY
    note_shape.line.color.rgb = LIGHT_GRAY
    note = slide.shapes.add_textbox(Inches(1.02), Inches(4.88), Inches(10.9), Inches(0.65))
    np = note.text_frame.paragraphs[0]
    np.text = "推荐组合：如果国内研发占主导，就先以 GLM-5.1 为试点主力；如果需要统一全球研发体验，就先以 GPT-5.4 为主力，再把 Claude Opus 4.6 接成升级通道。"
    np.font.name = FONT_FAMILY
    np.font.size = Pt(12)
    np.font.color.rgb = DARK
    _add_footer(slide, "目标不是一次选到最优模型，而是建立可持续的模型经营机制")


def _add_risk_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "风险与注意事项", "选型失败通常不是模型太弱，而是口径、路由和上下文管理失控")

    risks = [
        ("价格波动", "GLM-5.1 当前价格配置带有效期；后续可能调整。"),
        ("长上下文失控", "随着 agentic coding 深入，超长上下文比例会快速抬升。"),
        ("缓存收益不稳定", "团队工具链不同，缓存命中率可能与预期差异较大。"),
        ("跨币种比较偏差", "GLM 用人民币、OpenAI/Anthropic 用美元；采购决策需补充汇率和税费口径。"),
    ]
    for idx, (title, note) in enumerate(risks):
        top = Inches(1.85 + idx * 1.15)
        flag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), top, Inches(1.7), Inches(0.74))
        flag.fill.solid()
        flag.fill.fore_color.rgb = LIGHT_ORANGE
        flag.line.color.rgb = LIGHT_ORANGE
        fb = slide.shapes.add_textbox(Inches(0.98), top + Inches(0.16), Inches(1.3), Inches(0.25))
        fp = fb.text_frame.paragraphs[0]
        fp.text = title
        fp.font.name = FONT_FAMILY
        fp.font.size = Pt(12)
        fp.font.bold = True
        fp.font.color.rgb = NAVY

        nb = slide.shapes.add_textbox(Inches(2.78), top + Inches(0.13), Inches(9.25), Inches(0.45))
        np = nb.text_frame.paragraphs[0]
        np.text = note
        np.font.name = FONT_FAMILY
        np.font.size = Pt(12)
        np.font.color.rgb = DARK

    _add_footer(slide, "建议每月回顾一次：价格、上下文分布、缓存命中率、失败重试率")


def _add_appendix_slide(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _add_header(slide, "附录：来源与口径", "PPT 数字全部来自现有调研文档的公开口径")

    bullets = [
        "来源：OpenAI 官方 API Pricing 与 GPT-5.4 模型页。",
        "来源：Anthropic 官方 Claude Opus 4.6 页面、发布页与 rate limits 文档。",
        "来源：智谱官方模型概览、速率限制文档，以及官方价格页前端静态资源提取。",
        "用量模型假设：20 人、22 个工作日、三环节都深度使用大模型；不含 embedding、图像/音视频、税费和代理层加价。",
        "本 PPT 为管理汇报口径，不等同于企业协议价、私有部署价或真实采购合同价。",
    ]
    _add_bullets(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(2.5), bullets)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_GRAY
    box.line.color.rgb = LIGHT_GRAY
    note = slide.shapes.add_textbox(Inches(1.02), Inches(5.18), Inches(11.2), Inches(0.55))
    np = note.text_frame.paragraphs[0]
    np.text = "如需采购决策版，可进一步补充：人民币/美元统一汇率口径、税费口径、模型 SLA、数据合规要求、企业协议折扣。"
    np.font.name = FONT_FAMILY
    np.font.size = Pt(12)
    np.font.color.rgb = DARK

    _add_footer(slide, "源文档：2026-04-15-llm-pricing-and-usage-report.md")


if __name__ == "__main__":
    build_presentation(Path("2026-04-15-llm-pricing-and-usage-report.pptx"))
